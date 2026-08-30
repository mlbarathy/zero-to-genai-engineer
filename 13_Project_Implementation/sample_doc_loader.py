"""
Load RAG chunks from sample_docs/ in multiple file formats (docx, xlsx, pptx, pdf).
Used by build_db.py — mirrors M11 document-intelligence ingestion at teaching scale.
"""
from __future__ import annotations

import re
from pathlib import Path

DocRow = tuple[str, str, str, str, str, str, str]
# name, section, version, last_updated, chunk, source_type, source_file

VERSION_IN_NAME = re.compile(r"_v(\d+\.\d+)", re.I)
SECTION_VERSION = re.compile(r"^(.+?)\s*\(v(\d+\.\d+)\)\s*$", re.I)
DEFAULT_DATES = {
    "Discount_Policy": "2026-06-15",
    "Refund_and_Cancellation_Policy": "2026-07-01",
    "Employee_Handbook": "2026-05-20",
    "Food_Safety_SOP": "2026-04-10",
    "Menu_Copy_Style_Guide": "2026-06-01",
    "Opening_Hours_and_Reservations": "2026-03-01",
    "Inventory_Master": "2026-05-01",
    "Delivery_Operations_Manual": "2026-06-10",
    "Loyalty_Program_Rules": "2026-06-01",
    "Seating_and_Patio_Ops": "2026-05-15",
    "Payment_and_VAT_Policy": "2026-04-20",
    "Manager_Ops_Training": "2026-07-15",
    "Supplier_Price_List": "2026-06-20",
    "Municipality_Health_Inspection": "2026-05-10",
    "Promo_Calendar": "2026-06-15",
    "Delivery_Zones": "2026-06-10",
    "Marketing_AllHands": "2026-07-01",
}


def _doc_title_from_filename(path: Path) -> str:
    stem = VERSION_IN_NAME.sub("", path.stem)
    return stem.replace("_", " ").strip()


def _file_version(path: Path) -> str:
    m = VERSION_IN_NAME.search(path.stem)
    return f"v{m.group(1)}" if m else "v1.0"


def _last_updated(path: Path) -> str:
    stem = path.stem.split("_")[0]  # rough
    for key, date in DEFAULT_DATES.items():
        if path.stem.startswith(key) or key.replace("_", " ") in path.stem:
            return date
    return "2026-06-01"


def _parse_section_heading(text: str) -> tuple[str, str]:
    m = SECTION_VERSION.match(text.strip())
    if m:
        return m.group(1).strip(), f"v{m.group(2)}"
    return text.strip(), "v1.0"


def _sheet_to_text(rows: list[list]) -> str:
    lines = []
    for row in rows:
        cells = [str(c) if c is not None else "" for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def load_docx(path: Path) -> list[DocRow]:
    from docx import Document

    doc = Document(path)
    doc_name = _doc_title_from_filename(path)
    file_ver = _file_version(path)
    updated = _last_updated(path)
    rows: list[DocRow] = []
    section, sec_ver, buf = "Overview", file_ver, []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading 2" in style or (style.startswith("heading") and "2" in style):
            if buf:
                chunk = "\n".join(buf).strip()
                if chunk:
                    rows.append(
                        (doc_name, section, sec_ver, updated, chunk, "docx", path.name)
                    )
            section, sec_ver = _parse_section_heading(text)
            if sec_ver == "v1.0":
                sec_ver = file_ver
            buf = []
        elif "heading 1" in style:
            doc_name = text
        else:
            buf.append(text)
    if buf:
        chunk = "\n".join(buf).strip()
        if chunk:
            rows.append((doc_name, section, sec_ver, updated, chunk, "docx", path.name))
    return rows


def load_xlsx(path: Path) -> list[DocRow]:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    doc_name = _doc_title_from_filename(path)
    file_ver = _file_version(path)
    updated = _last_updated(path)
    rows: list[DocRow] = []
    for sheet in wb.worksheets:
        data = [list(row) for row in sheet.iter_rows(values_only=True)]
        chunk = _sheet_to_text(data)
        if chunk:
            rows.append(
                (
                    doc_name,
                    f"Sheet: {sheet.title}",
                    file_ver,
                    updated,
                    chunk,
                    "xlsx",
                    path.name,
                )
            )
    wb.close()
    return rows


def load_pptx(path: Path) -> list[DocRow]:
    from pptx import Presentation

    prs = Presentation(path)
    doc_name = _doc_title_from_filename(path)
    file_ver = _file_version(path)
    updated = _last_updated(path)
    rows: list[DocRow] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = f"Slide {i}"
        body_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape == slide.shapes.title:
                    title = shape.text.strip()
                else:
                    body_parts.append(shape.text.strip())
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        section, sec_ver = _parse_section_heading(title)
        if sec_ver == "v1.0":
            sec_ver = file_ver
        chunk = "\n".join(body_parts)
        if notes:
            chunk = (chunk + "\n\nSpeaker notes: " + notes).strip()
        if chunk:
            rows.append((doc_name, section, sec_ver, updated, chunk, "pptx", path.name))
    return rows


def load_pdf(path: Path) -> list[DocRow]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    doc_name = _doc_title_from_filename(path)
    file_ver = _file_version(path)
    updated = _last_updated(path)
    rows: list[DocRow] = []
    full = ""
    for page in reader.pages:
        full += (page.extract_text() or "") + "\n"
    # Split on lines that look like "Something (v1.0)"
    parts = re.split(r"\n(?=[A-Z][^\n]{3,40}\s*\(v\d+\.\d+\))", full)
    for part in parts:
        part = part.strip()
        if not part:
            continue
        lines = part.split("\n", 1)
        head = lines[0].strip()
        body = lines[1].strip() if len(lines) > 1 else part
        if SECTION_VERSION.match(head):
            section, sec_ver = _parse_section_heading(head)
            chunk = body
        else:
            section, sec_ver, chunk = "Content", file_ver, part
        if chunk:
            rows.append((doc_name, section, sec_ver, updated, chunk, "pdf", path.name))
    if not rows and full.strip():
        rows.append((doc_name, "Content", file_ver, updated, full.strip(), "pdf", path.name))
    return rows


def load_all_sample_docs(sample_dir: Path) -> list[DocRow]:
    """Walk sample_docs/{docx,xlsx,pptx,pdf}/ and return document rows."""
    if not sample_dir.is_dir():
        raise FileNotFoundError(f"Missing {sample_dir}")

    docs: list[DocRow] = []
    loaders = {
        ".docx": load_docx,
        ".xlsx": load_xlsx,
        ".pptx": load_pptx,
        ".pdf": load_pdf,
    }
    for ext, loader in loaders.items():
        folder = sample_dir / ext.lstrip(".")
        if not folder.is_dir():
            continue
        for path in sorted(folder.glob(f"*{ext}")):
            docs.extend(loader(path))

    if not docs:
        raise RuntimeError(
            f"No chunks loaded from {sample_dir}. Run: python generate_sample_docs.py"
        )
    return docs

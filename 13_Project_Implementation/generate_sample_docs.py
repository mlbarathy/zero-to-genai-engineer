#!/usr/bin/env python3
"""
Generate real sample_docs in docx / xlsx / pptx / pdf formats for the Dining Bot RAG corpus.
Run once (or after content edits):

    pip install python-docx openpyxl python-pptx fpdf2
    python generate_sample_docs.py
"""
from __future__ import annotations

from pathlib import Path

HERE = Path(__file__).resolve().parent
OUT = HERE / "sample_docs"


def _ensure_dirs() -> None:
    for sub in ("docx", "xlsx", "pptx", "pdf"):
        (OUT / sub).mkdir(parents=True, exist_ok=True)


def _write_docx_files() -> None:
    from docx import Document

    specs = [
        {
            "file": "Discount_Policy_v1.3.docx",
            "title": "Discount Policy",
            "sections": [
                (
                    "Manager limits (v1.2)",
                    "Restaurant managers may apply a discount of up to 15% on the order subtotal "
                    "without further approval. Discounts between 16% and 25% require shift supervisor "
                    "approval. Above 25% requires owner sign-off in writing. Log reason codes: "
                    "PROMO, COMPLAINT, STAFF_MEAL, PARTNER.",
                ),
                (
                    "Promotions (v1.2)",
                    "Standard promotional discounts are 5%, 10% and 15%. Promotions cannot be "
                    "combined with loyalty rewards on the same bill. Complimentary items are recorded "
                    "as comp lines, not discounts.",
                ),
                (
                    "Corporate and group bookings (v1.2)",
                    "Groups of 10+ may receive 12% off food subtotal when booked 48 hours ahead. "
                    "50% deposit required for groups above 20. Corporate invoice accounts use "
                    "negotiated rates — not ad-hoc POS discounts.",
                ),
            ],
        },
        {
            "file": "Refund_and_Cancellation_Policy_v2.1.docx",
            "title": "Refund Policy",
            "sections": [
                (
                    "Eligibility (v2.0)",
                    "Full refund when cancelled before preparation begins. After prep starts, "
                    "partial refund for undelivered items only. Quality complaints at manager "
                    "discretion up to affected item value.",
                ),
                (
                    "Method (v2.0)",
                    "Refund to original payment method. Card: 3–5 business days. Cash: immediate "
                    "from till with audit log. UPI/wallet within 24 hours via same terminal.",
                ),
                (
                    "Disputed card charges (v2.0)",
                    "Never duplicate-refund if guest opened a bank dispute. Escalate to owner with "
                    "order ID and payment reference. Chargeback response within five business days.",
                ),
            ],
        },
        {
            "file": "Employee_Handbook_v3.2.docx",
            "title": "Employee Handbook",
            "sections": [
                (
                    "Shift timings (v3.1)",
                    "Morning shift 10:00–16:00, evening 16:00–24:00. 30-minute break per shift, "
                    "11 hours minimum rest between shifts. Dining room opens 12:00.",
                ),
                (
                    "Leave (v3.1)",
                    "Two paid leave days accrued per month. Request seven days ahead except emergencies. "
                    "Max two kitchen staff on leave same day.",
                ),
                (
                    "Code of conduct (v3.2)",
                    "No phones on service floor except POS tablets. Staff meals in break room only. "
                    "Zero tolerance for intoxication on duty.",
                ),
            ],
        },
        {
            "file": "Opening_Hours_and_Reservations_v1.1.docx",
            "title": "Opening Hours Policy",
            "sections": [
                (
                    "Service hours (v1.0)",
                    "Open daily 12:00–23:00. Last kitchen orders 22:30. Closed 1 January each year.",
                ),
                (
                    "Ramadan hours (v1.1)",
                    "Ramadan dine-in 18:00–01:00. Limited iftar set 18:00–20:30.",
                ),
                (
                    "Private dining room (v1.1)",
                    "Private room seats 12; minimum spend AED 1,200 weekends. AV setup 24h notice.",
                ),
            ],
        },
        {
            "file": "Menu_Copy_Style_Guide_v1.2.docx",
            "title": "Menu Description Guide",
            "sections": [
                (
                    "Signature dishes (v1.1)",
                    "Butter Chicken: overnight-marinated tandoori chicken in tomato-butter gravy. "
                    "Chicken Biryani: Hyderabadi dum style. Always top-three on menus.",
                ),
                (
                    "Writing standards (v1.2)",
                    "Two sentences max per description. Use exact category names: Starters, Main Course, "
                    "Breads, Rice, Desserts, Beverages.",
                ),
                (
                    "Seasonal and inactive items (v1.1)",
                    "Set active=0 when discontinued — never delete history. Mushroom Pepper Fry and "
                    "Kulfi Falooda are inactive examples.",
                ),
            ],
        },
        {
            "file": "Payment_and_VAT_Policy_v1.1.docx",
            "title": "Payment Methods Policy",
            "sections": [
                (
                    "Accepted methods (v1.0)",
                    "Card, cash, UPI, wallet accepted. Card preferred above AED 200. Order counts as "
                    "revenue only when payment status is paid.",
                ),
                (
                    "Cash handling (v1.1)",
                    "Max AED 500 float at open. Safe drops every AED 1,500 with two witnesses.",
                ),
                (
                    "Split bills (v1.1)",
                    "Maximum three split payments per table. No IOUs.",
                ),
            ],
        },
    ]
    for spec in specs:
        doc = Document()
        doc.add_heading(spec["title"], level=1)
        for heading, body in spec["sections"]:
            doc.add_heading(heading, level=2)
            doc.add_paragraph(body)
        path = OUT / "docx" / spec["file"]
        doc.save(path)
        print("  docx:", path.name)


def _write_xlsx_files() -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    # Promo calendar
    wb = Workbook()
    ws = wb.active
    ws.title = "Weekday"
    ws.append(["day", "weekday", "promo_code", "discount_pct", "category", "time_start", "time_end"])
    for day in ("Mon", "Tue", "Wed", "Thu"):
        ws.append([day, "weekday", "WEEKDAY10", 10, "Main Course", "14:00", "17:00"])
    ws2 = wb.create_sheet("Blackout")
    ws2.append(["day", "note"])
    for day in ("Fri", "Sat", "Sun"):
        ws2.append([day, "No weekday lunch promo — peak / blackout"])
    wb.save(OUT / "xlsx" / "Promo_Calendar_2026_Q3.xlsx")
    print("  xlsx: Promo_Calendar_2026_Q3.xlsx")

    # Inventory
    wb = Workbook()
    ws = wb.active
    ws.title = "Reorder"
    ws.append(["ingredient", "unit", "stock", "reorder_level", "supplier"])
    rows = [
        ("Chicken", "kg", 24, 10, "Al Noor Meats"),
        ("Mutton", "kg", 9, 8, "Al Noor Meats"),
        ("Prawns", "kg", 5, 5, "Deira Fish Market"),
        ("Fish Fillet", "kg", 7, 8, "Deira Fish Market"),
        ("Basmati Rice", "kg", 60, 20, "Spice Trail Trading"),
        ("Garlic", "kg", 6, 3, "Spice Trail Trading"),
        ("Ginger", "kg", 5, 3, "Spice Trail Trading"),
    ]
    for r in rows:
        ws.append(list(r))
    wb.save(OUT / "xlsx" / "Inventory_Master_2026.xlsx")
    print("  xlsx: Inventory_Master_2026.xlsx")

    # Loyalty
    wb = Workbook()
    ws = wb.active
    ws.title = "Earning"
    ws.append(["rule", "value"])
    ws.append(["points_per_aed_food", "1 per AED 10 food (ex tax, delivery, tips)"])
    ws.append(["posting_delay_hours", 24])
    ws2 = wb.create_sheet("Redemption")
    ws2.append(["rule", "value"])
    ws2.append(["points_per_aed_off", "100 points = AED 10 off food subtotal"])
    ws2.append(["min_order_aed", 50])
    ws2.append(["stack_with_WEEKDAY10", "NO"])
    wb.save(OUT / "xlsx" / "Loyalty_Program_Rules_2026.xlsx")
    print("  xlsx: Loyalty_Program_Rules_2026.xlsx")

    # Supplier prices
    wb = Workbook()
    ws = wb.active
    ws.title = "Proteins"
    ws.append(["item", "unit_price_aed_per_kg", "vendor", "valid_until"])
    for row in [
        ("Chicken breast", 28.50, "Al Noor Meats", "2026-09-30"),
        ("Mutton leg", 42.00, "Al Noor Meats", "2026-09-30"),
        ("Fish fillet", 48.00, "Deira Fish Market", "2026-09-30"),
        ("Prawns 31/40", 65.00, "Deira Fish Market", "2026-09-30"),
    ]:
        ws.append(list(row))
    wb.save(OUT / "xlsx" / "Supplier_Price_List_Q3_2026.xlsx")
    print("  xlsx: Supplier_Price_List_Q3_2026.xlsx")

    # Delivery zones
    wb = Workbook()
    ws = wb.active
    ws.title = "Zones"
    ws.append(["zone_km", "fee_aed", "eta_min", "eta_max"])
    ws.append([5, 8, 35, 45])
    ws.append([8, 8, 40, 50])
    ws.append(["free_delivery_min_aed_weekday", 150, "", ""])
    wb.save(OUT / "xlsx" / "Delivery_Zones_2026.xlsx")
    print("  xlsx: Delivery_Zones_2026.xlsx")


def _write_pptx_files() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    def deck(filename: str, title: str, slides: list[tuple[str, str, str]]) -> None:
        prs = Presentation()
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Dining Bot sample corpus · GenAI-2026"
        for slide_title, body, notes in slides:
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide_title
            s.placeholders[1].text = body
            if s.has_notes_slide and notes:
                s.notes_slide.notes_text_frame.text = notes
        prs.save(OUT / "pptx" / filename)
        print("  pptx:", filename)

    deck(
        "Manager_Ops_Training_Aug2026.pptx",
        "Manager Ops Training",
        [
            (
                "Daily manager checklist (v1.0)",
                "Before lunch: cooler temps, low-stock report, 86 list, reservations, promo brief.\n"
                "Before dinner: repeat stock, patio per weather, delivery roster.",
                "Sign checklist in manager logbook.",
            ),
            (
                "Revenue discipline (v1.0)",
                "Paid revenue = SUM(orders.total) WHERE status='paid'.\n"
                "Never quote from memory — use analytics or dashboard.",
                "Finance cheat sheet on this slide for shift managers.",
            ),
            (
                "HITL and menu writes (v1.0)",
                "Only write path: ADD_MENU_ITEM after explicit human approval.\n"
                "Verify name, price, category, is_veg before approving.",
                "Matches Dining Bot capstone HITL demo.",
            ),
            (
                "Planning week ahead (v1.0)",
                "Combine: 7-day revenue SQL, low-stock xlsx, promo calendar, 3-day weather.\n"
                "Write weekly_plan.md with 3 manager actions — no DB writes from agent.",
                "Deep Agents planning harness follows this order.",
            ),
        ],
    )

    deck(
        "Seating_and_Patio_Ops_2026.pptx",
        "Seating & Patio Ops",
        [
            (
                "Indoor vs outdoor (v1.0)",
                "20 indoor tables, 8 patio tables.\n"
                "Patio Oct–Apr when highs below 38°C. Closed May–Sep unless owner event.",
                "Floor plan slide 3 is canonical table numbering.",
            ),
            (
                "Weather decisions (v1.1)",
                "Rain >40% or wind >35 km/h → close patio.\n"
                "32–38°C: fans/misting OK. No ad-hoc patio in summer.",
                "Use Open-Meteo forecast — not guest guesses.",
            ),
            (
                "Table allocation (v1.0)",
                "Parties 1–2: tables 1–8. 5+: combine 18–20 or patio 6–8 with notice.",
                "Wait incomplete large parties max 15 min in peak.",
            ),
        ],
    )

    deck(
        "Marketing_AllHands_July2026.pptx",
        "Marketing All-Hands",
        [
            (
                "Student and senior offers (v1.3)",
                "Students: 10% off food weekdays with valid ID.\n"
                "Seniors 60+: 15% Tuesday lunch only.\n"
                "Does NOT stack with WEEKDAY10 or loyalty.",
                "Owner discretion by location.",
            ),
            (
                "WEEKDAY10 campaign (v1.3)",
                "Mon–Thu 14:00–17:00: 10% off Main Course.\n"
                "POS code WEEKDAY10. Food subtotal only.",
                "Table tents and email subject approved.",
            ),
        ],
    )


def _write_pdf_files() -> None:
    from fpdf import FPDF

    def simple_pdf(filename: str, title: str, sections: list[tuple[str, str]]) -> None:
        pdf = FPDF()
        pdf.set_margins(18, 18, 18)
        pdf.set_auto_page_break(auto=True, margin=18)
        pdf.add_page()
        w = pdf.epw
        pdf.set_font("Helvetica", "B", 16)
        pdf.multi_cell(w, 10, title)
        pdf.ln(4)
        for heading, body in sections:
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w, 7, heading)
            pdf.set_font("Helvetica", size=11)
            safe = body.replace("\u2013", "-").replace("\u2014", "-").replace("\u00b0", " deg")
            pdf.multi_cell(w, 6, safe)
            pdf.ln(2)
        path = OUT / "pdf" / filename
        pdf.output(str(path))
        print("  pdf:", filename)

    simple_pdf(
        "Food_Safety_SOP_v1.1.pdf",
        "Food Safety SOP",
        [
            (
                "Storage (v1.0)",
                "Raw meat below 4 deg C, separate from veg and cooked food. FIFO rotation. "
                "Cooler log every 4 hours. Above 5 deg C for 30 min: notify manager.",
            ),
            (
                "Allergens (v1.0)",
                "Dairy, nuts, shellfish on menu. Prawn Koliwada and Fish Curry contain shellfish. "
                "Cannot guarantee nut-free kitchen — shared fryers.",
            ),
            (
                "Temperature and holding (v1.1)",
                "Hot-held curries above 63 deg C. Do not reheat rice twice. Danger zone 5-63 deg C "
                "over 2 hours: discard and log waste.",
            ),
        ],
    )

    simple_pdf(
        "Delivery_Operations_Manual_v1.1.pdf",
        "Delivery Operations Manual",
        [
            (
                "Order types (v1.0)",
                "dine_in, takeaway, delivery. Delivery fee AED 8 after tax — not in discount base.",
            ),
            (
                "Packaging standards (v1.0)",
                "Leak-proof containers. Hot items leave pass at 60 deg C+. Insulated bags over 20 min rides. "
                "Include allergen card.",
            ),
            (
                "Delivery refunds (v2.1)",
                "Full refund if rider has not left. Late >45 min: up to 20% partial or dessert — not both.",
            ),
        ],
    )

    simple_pdf(
        "Municipality_Health_Inspection_May2026.pdf",
        "Health Inspection Report",
        [
            (
                "Summary score (v1.0)",
                "Date 2026-05-08. Score 92/100. Grade A. Next visit expected September 2026.",
            ),
            (
                "Critical findings (v1.0)",
                "No critical violations. Grease trap log gap corrected on site. "
                "Allergen poster font size fixed within 48 hours.",
            ),
            (
                "Guest communication (v1.0)",
                "Approved social quote: Proud to maintain an A grade — thank you kitchen team. "
                "Do not claim 100% perfect.",
            ),
        ],
    )


def main() -> None:
    print("Generating sample_docs (docx, xlsx, pptx, pdf)…")
    _ensure_dirs()
    _write_docx_files()
    _write_xlsx_files()
    _write_pptx_files()
    _write_pdf_files()
    print("Done →", OUT)


if __name__ == "__main__":
    main()
